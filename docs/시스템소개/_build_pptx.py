# -*- coding: utf-8 -*-
"""SW 사업관리 시스템 개발 히스토리 — 편집 가능한 PPTX 생성"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 색 토큰 (정도UIT CI + 앱 Deep Teal) ----
NAVY  = RGBColor(0x2E, 0x31, 0x92)
TEAL  = RGBColor(0x0F, 0x7C, 0x8A)
TEALD = RGBColor(0x0A, 0x5B, 0x66)
TEALL = RGBColor(0xE3, 0xF2, 0xF4)
GREEN = RGBColor(0x58, 0xAA, 0x47)
GOLD  = RGBColor(0xD9, 0xA4, 0x41)
INK   = RGBColor(0x1F, 0x29, 0x33)
MUTED = RGBColor(0x7B, 0x87, 0x94)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE  = RGBColor(0xD8, 0xE0, 0xE5)
LIGHT = RGBColor(0xFA, 0xFC, 0xFC)

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
FONT = "Malgun Gothic"

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def grad(shape, c1, c2, angle=45):
    """간단한 선형 그라데이션 채우기 (좌상->우하)"""
    sp = shape.fill._xPr  # spPr
    # 기존 fill 제거
    for tag in ('a:noFill','a:solidFill','a:gradFill','a:blipFill','a:pattFill','a:grpFill'):
        for e in sp.findall(qn(tag)):
            sp.remove(e)
    grad = sp.makeelement(qn('a:gradFill'), {})
    lst = grad.makeelement(qn('a:gsLst'), {})
    for pos, col in ((0, c1), (100000, c2)):
        gs = grad.makeelement(qn('a:gs'), {'pos': str(pos)})
        clr = gs.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (col[0], col[1], col[2])})
        gs.append(clr); lst.append(gs)
    grad.append(lst)
    lin = grad.makeelement(qn('a:lin'), {'ang': str(angle*60000), 'scaled': '1'})
    grad.append(lin)
    # ln(테두리) 앞에 삽입
    ln = sp.find(qn('a:ln'))
    if ln is not None: ln.addprevious(grad)
    else: sp.append(grad)

def rect(s, x, y, w, h, fill=None, line=None, line_w=0.75, rounded=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    return shp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph = list of (text,size,color,bold)"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t
            r.font.name = FONT; r.font.size = Pt(sz)
            r.font.color.rgb = col; r.font.bold = bold
    return tb

def accent_bar(s):
    b = rect(s, 0, 0, 0.42, 7.5, fill=NAVY)
    grad(b, (0x2E,0x31,0x92), (0x0F,0x7C,0x8A), angle=90)

def footer(s, n):
    txt(s, 0.95, 7.05, 4, 0.3, [[("SYSTEM INTRO", 9, MUTED, False)]])
    txt(s, 12.0, 7.05, 1.0, 0.3, [[("%02d" % n, 10, MUTED, False)]], align=PP_ALIGN.RIGHT)

def title_block(s, eyebrow, title, sub=None):
    txt(s, 0.95, 0.55, 11.5, 0.4, [[(eyebrow, 13, TEAL, True)]])
    txt(s, 0.93, 0.95, 11.6, 0.9, [[(title, 30, NAVY, True)]], line_spacing=1.05)
    if sub:
        txt(s, 0.95, 1.75, 11.5, 0.5, [[(sub, 15, MUTED, False)]])

# ============================================================ 1 표지
s = slide(); bg(s, NAVY)
cover = rect(s, 0, 0, 13.333, 7.5, fill=NAVY); grad(cover, (0x2E,0x31,0x92), (0x0F,0x7C,0x8A), 60)
txt(s, 0.95, 1.6, 11, 0.4, [[("SW PROJECT MANAGEMENT SYSTEM", 14, RGBColor(0xCF,0xD8,0xEC), True)]])
txt(s, 0.92, 2.05, 11.5, 1.8, [
    [("SW 사업관리 시스템", 50, WHITE, True)],
    [("개발 히스토리", 50, WHITE, True)],
], line_spacing=1.05)
txt(s, 0.95, 3.95, 11, 0.5, [[("처음 한 줄부터 지금까지 — 3개월의 기록", 20, RGBColor(0xE6,0xEE,0xF2), False)]])
nums = [("335","작업 기록"),("41일","작업한 날"),("≈181h","추정 작업시간"),("4단계","성장 과정")]
for i,(n,l) in enumerate(nums):
    x = 0.95 + i*2.45
    txt(s, x, 4.85, 2.3, 0.7, [[(n, 40, WHITE, True)]])
    txt(s, x, 5.65, 2.3, 0.4, [[(l, 13, RGBColor(0xCF,0xD8,0xEC), False)]])
txt(s, 0.95, 6.85, 11, 0.4, [[("2026-04-03 ~ 2026-06-17  ·  박욱진(SW지원부)  ·  2026-06-19", 13, RGBColor(0xD3,0xDC,0xE6), False)]])

# ============================================================ 2 WHAT
s = slide(); bg(s, WHITE); accent_bar(s)
title_block(s, "WHAT", "한 화면에서 보는 우리 부서의 “관제탑”",
            "흩어져 있던 사업·서류·점검·장애 대응을 하나의 웹 시스템으로")
tiles = [("📊","대시보드","사업 현황·금액 요약"),("📁","사업 관리","등록·검색·단계 관리"),
         ("📄","문서 자동화","착수·기성·준공계"),("🔍","점검내역서","현장 점검→PDF"),
         ("💰","견적서","인건비 자동 계산"),("🛠","장애·업무지원","대응 내역 기록"),
         ("🖥","인프라 자산","서버·장비 관리"),("👤","권한·로그","승인·기록·통계")]
tw, th, gx, gy = 2.78, 1.65, 0.18, 0.22
x0, y0 = 0.95, 2.45
for i,(ic,tt,td) in enumerate(tiles):
    r,c = divmod(i,4)
    x = x0 + c*(tw+gx); y = y0 + r*(th+gy)
    rect(s, x, y, tw, th, fill=LIGHT, line=LINE, rounded=True)
    txt(s, x+0.22, y+0.16, tw-0.4, 0.5, [[(ic, 24, INK, False)]])
    txt(s, x+0.22, y+0.72, tw-0.4, 0.4, [[(tt, 16, TEALD, True)]])
    txt(s, x+0.22, y+1.10, tw-0.4, 0.4, [[(td, 12, MUTED, False)]])
footer(s, 2)

# ============================================================ 3 숫자
s = slide(); bg(s, WHITE); accent_bar(s)
title_block(s, "BY THE NUMBERS", "숫자로 보는 3개월", "2026년 4월 3일 첫 저장 → 6월 17일까지")
cards = [("335","작업 기록(커밋)","= 작업 일지 한 칸씩",(NAVY,TEALD)),
         ("41일","실제 작업한 날","전체 76일 중",(TEALD,TEAL)),
         ("181h","추정 작업시간","하루 평균 4.4h",(TEAL,RGBColor(0x1A,0x9A,0xAA))),
         ("70+","화면(페이지)","처리 모듈 20+",(GREEN,TEAL))]
cw, gx = 2.86, 0.22; x0, y0 = 0.95, 2.55
for i,(n,l,d,(c1,c2)) in enumerate(cards):
    x = x0 + i*(cw+gx)
    bxshp = rect(s, x, y0, cw, 2.35, fill=c1, rounded=True)
    grad(bxshp, (c1[0],c1[1],c1[2]), (c2[0],c2[1],c2[2]), 60)
    txt(s, x+0.28, y0+0.3, cw-0.5, 0.9, [[(n, 46, WHITE, True)]])
    txt(s, x+0.28, y0+1.35, cw-0.5, 0.5, [[(l, 15, WHITE, True)]])
    txt(s, x+0.28, y0+1.78, cw-0.5, 0.4, [[(d, 12, RGBColor(0xE6,0xEE,0xF2), False)]])
rect(s, 0.95, 5.5, 11.42, 0.95, fill=TEALL, rounded=True)
txt(s, 1.25, 5.62, 10.9, 0.75,
    [[("💡  화면 수에 비해 ", 15, INK, False),("보안 점검·데이터 정리·문서화", 15, TEALD, True),
      ("에 들인 비중이 큽니다 — “빠르게”보다 “틀리지 않게”에 무게.", 15, INK, False)]],
    anchor=MSO_ANCHOR.MIDDLE)
footer(s, 3)

# ============================================================ 4 4단계
s = slide(); bg(s, WHITE); accent_bar(s)
title_block(s, "JOURNEY", "시스템이 자라온 4단계")
steps = [("1단계","4월 초·중순","기능을 빠르게 채우다","문서·견적·점검내역서 확장, 태블릿 대응, 팀 모니터링"),
         ("2단계","4월 중·하순","점검하고 튼튼하게 정리","보안 감사, 데이터 정규화, 개발 문서 전면 재편"),
         ("3단계","5월","현장 점검 자동화 + 디자인 통일","QR·원격 점검 프로그램, 40여 화면 색 통일"),
         ("4단계","6월","실무에 맞춰 넓히다","업무지원·장애처리, 지식베이스, 사이드바·HTTPS")]
cw, gx = 2.78, 0.18; x0, y0 = 0.95, 2.2
for i,(ph,when,st,sd) in enumerate(steps):
    x = x0 + i*(cw+gx)
    pill = rect(s, x, y0, 1.15, 0.42, fill=TEALL, rounded=True)
    txt(s, x, y0+0.04, 1.15, 0.34, [[(ph, 13, TEALD, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x, y0+0.55, cw, 0.35, [[(when, 13, MUTED, False)]])
    txt(s, x, y0+0.95, cw, 1.1, [[(st, 18, NAVY, True)]], line_spacing=1.05)
    txt(s, x, y0+2.05, cw, 1.2, [[(sd, 13, RGBColor(0x55,0x63,0x6E), False)]], line_spacing=1.2)
    if i < 3:
        txt(s, x+cw-0.02, y0+0.75, 0.3, 0.4, [[("→", 22, TEAL, True)]])
rect(s, 0.95, 5.7, 11.42, 0.8, fill=LIGHT, line=LINE, rounded=True)
txt(s, 1.25, 5.82, 10.9, 0.6,
    [[("집을 짓듯 — ", 15, INK, False),
      ("① 방 늘리기 → ② 배관·전기 정비 → ③ 인테리어 통일 → ④ 가구 채우기", 15, TEALD, True)]],
    anchor=MSO_ANCHOR.MIDDLE)
footer(s, 4)

# ============================================================ 5 월별
s = slide(); bg(s, WHITE); accent_bar(s)
title_block(s, "WORKLOAD · MONTHLY", "월별 작업량 & 작업시간")
# 좌: 막대
bars = [("4월",65.5,1.00),("5월",60.7,0.927),("6월",54.6,0.834)]
bx, by, bw = 0.95, 2.6, 5.0
for i,(lbl,val,ratio) in enumerate(bars):
    y = by + i*0.85
    txt(s, bx, y, 0.7, 0.45, [[(lbl, 16, MUTED, False)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, bx+0.85, y, bw, 0.5, fill=RGBColor(0xEE,0xF2,0xF4), rounded=True)
    fillw = bw*ratio
    fb = rect(s, bx+0.85, y, fillw, 0.5, fill=TEAL, rounded=True)
    grad(fb, (0x0F,0x7C,0x8A), (0x2E,0x31,0x92), 0)
    txt(s, bx+0.85, y, fillw-0.15, 0.5, [[("%.1fh"%val, 14, WHITE, True)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
# 우: 표
tx, ty = 6.7, 2.55
rows = [("월","기록","일수","시간","평균"),
        ("4월","104","13","65.5h","5.0h"),
        ("5월","138","18","60.7h","3.4h"),
        ("6월","93","10","54.6h","5.5h"),
        ("합계","335","41","180.8h","4.4h")]
tbl = s.shapes.add_table(5, 5, Inches(tx), Inches(ty), Inches(5.6), Inches(2.7)).table
widths = [1.1,1.05,1.05,1.25,1.15]
for ci,wv in enumerate(widths): tbl.columns[ci].width = Inches(wv)
for ri,row in enumerate(rows):
    tbl.rows[ri].height = Inches(0.54)
    for ci,val in enumerate(row):
        cell = tbl.cell(ri,ci)
        cell.margin_left = cell.margin_right = Pt(6)
        cell.margin_top = cell.margin_bottom = Pt(2)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if ci==0 else PP_ALIGN.RIGHT
        r = p.add_run(); r.text = val; r.font.name = FONT; r.font.size = Pt(13)
        if ri==0:
            r.font.bold=True; r.font.color.rgb=WHITE; cell.fill.solid(); cell.fill.fore_color.rgb=NAVY
        elif ri==4:
            r.font.bold=True; r.font.color.rgb=TEALD; cell.fill.solid(); cell.fill.fore_color.rgb=TEALL
        else:
            r.font.color.rgb=INK; cell.fill.solid(); cell.fill.fore_color.rgb=(LIGHT if ri%2 else WHITE)
txt(s, 0.95, 6.45, 11.4, 0.5,
    [[("작업시간은 저장 기록(커밋) 시각 기반 ",13,MUTED,False),("추정치",13,TEALD,True),
      (" — 간격 2시간 미만=연속작업, 묶음마다 +1h.",13,MUTED,False)]])
footer(s, 5)

# ============================================================ 6 일별 차트
s = slide(); bg(s, WHITE); accent_bar(s)
title_block(s, "WORKLOAD · DAILY", "일별 작업시간 (41일)")
heights = [7,7,30,44,31,69,38,7,43,41,57,46,32,14,15,7,7,28,14,7,17,35,69,7,7,32,7,7,19,28,100,34,63,7,52,33,63,43,11,44,28]
cx, cy, cw, ch = 0.95, 2.4, 11.4, 3.4
n = len(heights); gap = 0.04
colw = (cw - gap*(n-1)) / n
base = cy + ch
for i,hh in enumerate(heights):
    x = cx + i*(colw+gap)
    bh = ch * (hh/100.0)
    col = GOLD if hh==100 else TEAL
    cbar = rect(s, x, base-bh, colw, bh, fill=col, rounded=False)
    if hh==100:
        grad(cbar,(0xD9,0xA4,0x41),(0xC0,0x53,0x2B),90)
    else:
        grad(cbar,(0x0F,0x7C,0x8A),(0x2E,0x31,0x92),90)
rect(s, cx, base, cw, 0.02, fill=LINE)
for lbl,frac in [("4/3",0),("4/26",0.27),("5/17",0.55),("5/31",0.73),("6/11",0.9),("6/17",1.0)]:
    txt(s, cx+cw*frac-0.4, base+0.08, 0.8, 0.3, [[(lbl,11,MUTED,False)]], align=PP_ALIGN.CENTER)
txt(s, 0.95, 6.35, 11.4, 0.7,
    [[("최장일 ",14,RGBColor(0x55,0x63,0x6E),False),("5/31 ≈14.5h(52건)",14,TEALD,True),
      (" — 화면 40여 개 디자인 일괄 통일. 다음 ",14,RGBColor(0x55,0x63,0x6E),False),
      ("4/19·5/17 ≈10h",14,TEALD,True),(". 주말·심야 작업이 많았습니다.",14,RGBColor(0x55,0x63,0x6E),False)]],
    line_spacing=1.2)
footer(s, 6)

# ============================================================ 7~9 단계별 패널
def panel_slide(no, eyebrow, title, left, right, quote):
    s = slide(); bg(s, WHITE); accent_bar(s)
    title_block(s, eyebrow, title)
    px, py, pw, ph = 0.95, 2.05, 5.55, 3.3
    for idx,(head,items) in enumerate([left,right]):
        x = px + idx*(pw+0.32)
        rect(s, x, py, pw, ph, fill=LIGHT, line=LINE, rounded=True)
        txt(s, x+0.35, py+0.28, pw-0.7, 0.5, [[(head, 19, TEALD, True)]])
        runs=[]
        for it in items:
            if isinstance(it, tuple):
                runs.append([("• ",16,TEAL,True),(it[0],16,INK,True),("  "+it[1],13,MUTED,False)])
            else:
                runs.append([("• ",16,TEAL,True),(it,16,INK,False)])
        txt(s, x+0.35, py+0.95, pw-0.7, ph-1.1, runs, space_after=9, line_spacing=1.1)
    if quote:
        rect(s, 0.95, 5.7, 0.07, 0.85, fill=TEAL)
        txt(s, 1.2, 5.72, 11.0, 0.85, [[(quote, 19, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    footer(s, no)

panel_slide(7, "DEEP DIVE · APR", "4월 — 채우고, 다지다",
    ("① 기능을 빠르게", [("문서·견적·점검","대규모 확장"),("태블릿 대응","현장 갤럭시탭 세로"),"팀 진행율 대시보드"]),
    ("② 점검하고 정리", [("보안 감사","codex 기반 조치"),("데이터 정규화","중복 제거·표준화"),"개발 문서 전면 재편"]),
    "“기능을 늘린 뒤, 바로 안전성과 구조를 점검했다.”")

panel_slide(8, "DEEP DIVE · MAY", "5월 — 자동화와 통일",
    ("🤖 현장 점검 자동화", [("QR 반출","폐쇄망 결과 반출"),("원격 점검","Telnet·SSH, 구형 서버 호환"),("점검표 자동 생성","수집값→양식+그래프")]),
    ("🎨 디자인 전면 통일", [("색 기준표(토큰)","도입"),("40여 화면 일괄 정비","5/31 하루 52건"),"브랜드 색으로 통일"]),
    "“손이 많이 가던 현장 점검을, 현장에 두는 작은 프로그램이 대신하게.”")

# 9는 단일 리스트형
s = slide(); bg(s, WHITE); accent_bar(s)
title_block(s, "DEEP DIVE · JUN", "6월 — 실무에 맞춰 넓히다")
items9 = [
    ("서류 편의","날인본 스캔 보관 · 공문 서식 보정 · 검색결과 일괄 ZIP 다운로드"),
    ("화면 구조 개편","상단 메뉴를 왼쪽 사이드바로 + 모바일 반응형, 대시보드를 메인으로"),
    ("업무계획 캐스케이드","시·도 → 시·군·구 → 시스템으로 이어 선택"),
    ("장애처리·업무지원 개편","직원 85명 조직도 · 엔지니어/요청자 · 지식베이스 추천(증상→원인→조치)"),
    ("로그·운영문서 정리","통계 추가, 상세·수정·삭제 흐름 정돈"),
]
runs=[]
for h,d in items9:
    runs.append([("•  ",18,TEAL,True),(h+" — ",18,TEALD,True),(d,16,RGBColor(0x33,0x42,0x4E),False)])
txt(s, 0.95, 2.2, 11.5, 3.4, runs, space_after=13, line_spacing=1.1)
rect(s, 0.95, 5.95, 0.07, 0.8, fill=TEAL)
txt(s, 1.2, 5.97, 11.0, 0.8, [[("“실제로 살아보며 불편을 없애는 단계 — 가구를 채워 넣다.”", 19, NAVY, True)]],
    anchor=MSO_ANCHOR.MIDDLE)
footer(s, 9)

# ============================================================ 10 일하는 방식
s = slide(); bg(s, WHITE); accent_bar(s)
title_block(s, "HOW WE WORK", "일하는 방식 — 검토를 거쳐 안전하게",
            "바로 만들지 않고, 계획→검토→승인→구현→검증을 거쳤습니다")
flow = [("요청",False),("기획서",False),("AI(codex) 검토",False),("사람 승인",True),
        ("구현",False),("AI 검증",False),("완료",True)]
fx, fy = 0.95, 2.35;
xx = fx
for i,(label,key) in enumerate(flow):
    w = 0.55 + len(label)*0.135
    nb = rect(s, xx, fy, w, 0.62, fill=(NAVY if key else TEALL), line=(NAVY if key else RGBColor(0xC5,0xDD,0xE0)), rounded=True)
    txt(s, xx, fy+0.02, w, 0.58, [[(label, 14, (WHITE if key else TEALD), True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    xx += w
    if i < len(flow)-1:
        txt(s, xx, fy+0.05, 0.42, 0.55, [[("→", 18, TEAL, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        xx += 0.42
# 하단 2패널
py, pw, ph = 3.55, 5.55, 2.3
panels=[("가상 팀 체계",[("기획·DB·디자인·개발·테스트","5역할 분담 점검"),("화면 변경 시","디자인 자문 필수")]),
        ("“안 하기로 한 결정”도 기록",[("계약서 PDF 파싱","만들었다가 취소"),("일부 보안 작업","의도적 보류")])]
for idx,(head,items) in enumerate(panels):
    x = 0.95 + idx*(pw+0.32)
    rect(s, x, py, pw, ph, fill=LIGHT, line=LINE, rounded=True)
    txt(s, x+0.35, py+0.25, pw-0.7, 0.5, [[(head,18,TEALD,True)]])
    runs=[[("• ",15,TEAL,True),(a,15,INK,True),("  "+b,13,MUTED,False)] for a,b in items]
    txt(s, x+0.35, py+0.9, pw-0.7, ph-1.0, runs, space_after=10, line_spacing=1.1)
footer(s, 10)

# ============================================================ 11 맺음말
s = slide(); bg(s, TEALD)
cover = rect(s, 0, 0, 13.333, 7.5, fill=TEALD); grad(cover,(0x0A,0x5B,0x66),(0x2E,0x31,0x92),60)
txt(s, 0.95, 1.5, 11, 0.4, [[("SUMMARY", 13, RGBColor(0xCF,0xD8,0xEC), True)]])
txt(s, 0.92, 2.0, 11.6, 1.6, [
    [("단순 사업 목록에서", 38, WHITE, True)],
    [("부서의 통합 업무 플랫폼으로", 38, WHITE, True)],
], line_spacing=1.1)
txt(s, 0.95, 3.95, 11.3, 1.0, [
    [("41일 · 약 181시간 · 335번의 저장이 쌓여 —", 18, RGBColor(0xE6,0xEE,0xF2), False)],
    [("문서 자동화 · 점검 자동화 · 견적 · 장애/업무지원 · 지식베이스 · 통일된 디자인을 갖췄습니다.", 18, RGBColor(0xE6,0xEE,0xF2), False)],
], line_spacing=1.3)
txt(s, 0.95, 5.4, 11.3, 0.9, [
    [("무엇보다, 모든 과정이 ", 16, RGBColor(0xCF,0xD8,0xEC), False),
     ("계획→검토→승인→구현→검증", 16, WHITE, True),
     ("으로 기록되어", 16, RGBColor(0xCF,0xD8,0xEC), False)],
    [("앞으로도 안전하게 이어갈 수 있습니다.", 16, RGBColor(0xCF,0xD8,0xEC), False)],
], line_spacing=1.3)
txt(s, 0.95, 6.85, 11, 0.4, [[("SW 사업관리 시스템 · 개발 히스토리 · 2026-06-19", 13, RGBColor(0xD3,0xDC,0xE6), False)]])

out = r"C:\Users\PUJ\eclipse-workspace\swmanager\docs\시스템소개\SW사업관리시스템_개발히스토리.pptx"
prs.save(out)
print("saved:", out)
print("slides:", len(prs.slides._sldIdLst))
